import os
import io
import zipfile
import logging
import subprocess
from flask import Flask, request, send_file, render_template, send_from_directory
from werkzeug.utils import secure_filename
from PIL import Image
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pdf2docx import Converter

# Initialize Flask
app = Flask(__name__)

UPLOAD_FOLDER = "uploads"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

CONVERTED_FOLDER = "converted"
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50 MB

logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)


def convert_with_libreoffice(input_path, output_dir):
    """
    Converts a document to PDF using LibreOffice with high fidelity.
    Returns the path to the converted file or None on failure.
    """
    try:
        command = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            input_path,
        ]
        logger.info(f"Running LibreOffice command: {' '.join(command)}")
        process = subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120  # 2-minute timeout for large files
        )

        if process.returncode != 0:
            error_message = process.stderr.decode("utf-8", errors="ignore")
            logger.error(f"LibreOffice conversion failed: {error_message}")
            return None

        original_filename = os.path.basename(input_path)
        pdf_filename = os.path.splitext(original_filename)[0] + ".pdf"
        output_path = os.path.join(output_dir, pdf_filename)

        if os.path.exists(output_path):
            logger.info(f"Successfully converted file to {output_path}")
            return output_path
        return None
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        logger.error(f"LibreOffice execution error: {e}. Is LibreOffice installed and in PATH?")
        return None


@app.route("/")
def home():
    return render_template("index.html")


@app.route("/convert", methods=["POST", "OPTIONS"])
def convert_file():
    if request.method == "OPTIONS":
        return '', 200
    
    if "file" not in request.files:
        return "No file uploaded.", 400

    uploaded_file = request.files["file"]
    target_format = request.form.get("target")

    if not uploaded_file:
        return "No file selected.", 400

    if not target_format:
        return "Target format not selected.", 400

    filename = secure_filename(uploaded_file.filename)
    ext = os.path.splitext(filename)[1].lower()
    name = os.path.splitext(filename)[0]
    target_format = target_format.lower()

    temp_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    uploaded_file.save(temp_path)

    try:
        # ==============================
        # IMAGE → PNG/JPG/PDF/DOCX/PPTX
        # ==============================
        if ext in [".jpg", ".jpeg", ".png"]:
            img = Image.open(temp_path).convert("RGB")

            if target_format in ["jpg", "jpeg", "png"]:
                buffer = io.BytesIO()
                fmt = "JPEG" if target_format in ["jpg", "jpeg"] else target_format.upper()
                img.save(buffer, fmt)
                buffer.seek(0)
                return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")

            elif target_format == "pdf":
                buffer = io.BytesIO()
                img.save(buffer, "PDF", resolution=100.0)
                buffer.seek(0)
                return send_file(buffer, as_attachment=True, download_name=f"{name}.pdf")

            elif target_format in ["docx", "doc"]:
                doc = Document()
                doc.add_picture(temp_path)
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")

            elif target_format in ["pptx", "ppt"]:
                ppt = Presentation()
                slide = ppt.slides.add_slide(ppt.slide_layouts[6])
                slide.shapes.add_picture(temp_path, 0, 0, width=ppt.slide_width)
                buffer = io.BytesIO()
                ppt.save(buffer)
                buffer.seek(0)
                return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")

        # ==============================
        # PDF → DOCX/PPTX/PNG/JPG
        # ==============================
        elif ext == ".pdf":
            if target_format in ["png", "jpg", "jpeg"]:
                logger.info(f"PDF → {target_format}")
                images = convert_from_path(temp_path)

                if len(images) == 1:
                    fmt = "PNG" if target_format == "png" else "JPEG"
                    buffer = io.BytesIO()
                    images[0].save(buffer, format=fmt)
                    buffer.seek(0)
                    return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")

                else:
                    zip_buffer = io.BytesIO()
                    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
                        for i, img in enumerate(images):
                            img_buffer = io.BytesIO()
                            fmt = "PNG" if target_format == "png" else "JPEG"
                            img.save(img_buffer, format=fmt)
                            img_buffer.seek(0)
                            zip_file.writestr(f"{name}_page_{i+1}.{target_format}", img_buffer.getvalue())
                    zip_buffer.seek(0)
                    return send_file(zip_buffer, as_attachment=True, download_name=f"{name}_{target_format}s.zip", mimetype="application/zip")

            elif target_format in ["docx", "doc"]:
                try:
                    buffer = io.BytesIO()
                    cv = Converter(temp_path)
                    cv.convert(buffer, start=0, end=None)
                    cv.close()
                    buffer.seek(0)
                    return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")
                except Exception as e:
                    logger.error(f"PDF → DOCX error: {e}")
                    return f"Error converting PDF to Word: {str(e)}", 500

            elif target_format in ["pptx", "ppt"]:
                try:
                    images = convert_from_path(temp_path)
                    prs = Presentation()
                    for img in images:
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        img_buffer = io.BytesIO()
                        img.save(img_buffer, "PNG")
                        img_buffer.seek(0)
                        slide.shapes.add_picture(img_buffer, 0, 0, width=prs.slide_width)
                    
                    buffer = io.BytesIO()
                    prs.save(buffer)
                    buffer.seek(0)
                    return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")
                except Exception as e:
                    logger.error(f"PDF → PPTX error: {e}")
                    return f"Error converting PDF to PowerPoint: {str(e)}", 500

        # ==============================
        # WORD → PDF/PPT
        # ==============================
        elif ext in [".doc", ".docx"]:
            try:
                doc = Document(temp_path)
            except Exception as e:
                return f"Error reading Word file: {str(e)}", 400

            if target_format == "pdf":
                # --- HIGH-FIDELITY CONVERSION ---
                output_pdf_path = convert_with_libreoffice(temp_path, CONVERTED_FOLDER)
                if output_pdf_path:
                    return send_from_directory(CONVERTED_FOLDER, os.path.basename(output_pdf_path), as_attachment=True)
                else:
                    return "Error during DOCX to PDF conversion.", 500

            elif target_format in ["pptx", "ppt"]:
                ppt = Presentation()
                for para in doc.paragraphs:
                    if para.text.strip():
                        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                        slide.shapes.title.text = para.text[:60]
                
                buffer = io.BytesIO()
                ppt.save(buffer)
                buffer.seek(0)
                return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")

        # ==============================
        # PPT → PDF/WORD
        # ==============================
        elif ext in [".ppt", ".pptx"]:
            ppt = Presentation(temp_path)

            if target_format == "pdf":
                # --- HIGH-FIDELITY CONVERSION ---
                output_pdf_path = convert_with_libreoffice(temp_path, CONVERTED_FOLDER)
                if output_pdf_path:
                    return send_from_directory(CONVERTED_FOLDER, os.path.basename(output_pdf_path), as_attachment=True)
                else:
                    return "Error during PPTX to PDF conversion.", 500

            elif target_format in ["doc", "docx"]:
                doc = Document()
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            doc.add_paragraph(shape.text)
                
                buffer = io.BytesIO()
                doc.save(buffer)
                buffer.seek(0)
                return send_file(buffer, as_attachment=True, download_name=f"{name}.{target_format}")

        # ==============================
        # Unsupported
        # ==============================
        return f"Unsupported conversion: {ext} → {target_format}", 400

    except Exception as e:
        logger.error(f"Conversion error: {e}", exc_info=True)
        return f"Conversion error: {str(e)}", 500

    finally:
        if os.path.exists(temp_path):
            try:
                # Clean up original upload
                os.remove(temp_path)
                # Clean up converted file if it exists
                # Note: This is a simple cleanup. A better approach for production would be a scheduled task.
            except Exception as e:
                pass


if __name__ == "__main__":
    app.run(debug=True)
